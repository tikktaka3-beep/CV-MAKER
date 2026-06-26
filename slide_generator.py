"""
slide_generator.py — Ko'p theme'li dizayn
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ==================== THEME DEFINITIONS ====================
THEMES = {
    "rasmiy": {
        "name": "Rasmiy", "bg_dark": RGBColor(0x1E, 0x27, 0x61), "accent": RGBColor(0xF2, 0xA9, 0x3B),
        "bg_light": RGBColor(0xF8, 0xF9, 0xFA), "text": RGBColor(0x2B, 0x2B, 0x2B)
    },
    "zamonaviy": {
        "name": "Zamonaviy", "bg_dark": RGBColor(0x11, 0x11, 0x11), "accent": RGBColor(0x00, 0xFF, 0xAA),
        "bg_light": RGBColor(0xF0, 0xF0, 0xF0), "text": RGBColor(0x22, 0x22, 0x22)
    },
    "rangbarang": {
        "name": "Rang-barang", "bg_dark": RGBColor(0x2A, 0x1B, 0x4A), "accent": RGBColor(0xFF, 0x6B, 0x6B),
        "bg_light": RGBColor(0xFF, 0xF9, 0xE6), "text": RGBColor(0x2C, 0x3E, 0x50)
    },
    "och": {
        "name": "Och ranglar", "bg_dark": RGBColor(0xE3, 0xF2, 0xFD), "accent": RGBColor(0x29, 0xB6, 0xF6),
        "bg_light": RGBColor(0xFA, 0xFB, 0xFF), "text": RGBColor(0x15, 0x65, 0xC0)
    },
    "toq": {
        "name": "To'q ranglar", "bg_dark": RGBColor(0x0A, 0x1F, 0x3D), "accent": RGBColor(0xFF, 0xC1, 0x07),
        "bg_light": RGBColor(0xF4, 0xF4, 0xF9), "text": RGBColor(0x1B, 0x26, 0x3B)
    },
    "qadimiy": {
        "name": "Qadimiy", "bg_dark": RGBColor(0x3C, 0x2F, 0x2F), "accent": RGBColor(0xD4, 0xAF, 0x37),
        "bg_light": RGBColor(0xF5, 0xF0, 0xE6), "text": RGBColor(0x3C, 0x2F, 0x2F)
    },
    "adabiy": {
        "name": "Adabiy", "bg_dark": RGBColor(0x2C, 0x18, 0x1F), "accent": RGBColor(0xE8, 0xB9, 0xAB),
        "bg_light": RGBColor(0xFA, 0xF4, 0xED), "text": RGBColor(0x3E, 0x27, 0x2A)
    },
    "matematik": {
        "name": "Matematik", "bg_dark": RGBColor(0x1A, 0x2A, 0x44), "accent": RGBColor(0x00, 0xFF, 0xFF),
        "bg_light": RGBColor(0xF0, 0xF8, 0xFF), "text": RGBColor(0x0F, 0x2B, 0x4A)
    },
    "kimyoviy": {
        "name": "Kimyoviy", "bg_dark": RGBColor(0x2E, 0x1A, 0x3A), "accent": RGBColor(0xFF, 0x00, 0x7F),
        "bg_light": RGBColor(0xF8, 0xF0, 0xFF), "text": RGBColor(0x4A, 0x23, 0x5A)
    },
    "fizik": {
        "name": "Fizik", "bg_dark": RGBColor(0x0F, 0x2B, 0x4A), "accent": RGBColor(0x39, 0xFF, 0x14),
        "bg_light": RGBColor(0xE6, 0xF3, 0xFF), "text": RGBColor(0x1C, 0x3A, 0x5E)
    },
    "astronomik": {
        "name": "Astronomik", "bg_dark": RGBColor(0x0B, 0x0C, 0x2E), "accent": RGBColor(0xAA, 0xFF, 0xFF),
        "bg_light": RGBColor(0xE0, 0xF0, 0xFF), "text": RGBColor(0x1E, 0x2A, 0x5E)
    },
    "sport": {
        "name": "Sport", "bg_dark": RGBColor(0x1B, 0x2A, 0x1F), "accent": RGBColor(0xFF, 0x2D, 0x00),
        "bg_light": RGBColor(0xF0, 0xF7, 0xF0), "text": RGBColor(0x1B, 0x2F, 0x1F)
    }
}

def get_theme(theme_key: str):
    return THEMES.get(theme_key, THEMES["rasmiy"])

# ==================== YORDAMCHI FUNKSIYALAR ====================
def _set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_text(slide, left, top, width, height, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if anchor: tf.vertical_anchor = anchor
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri Light" if "zamonaviy" in theme_key else "Calibri"
    return box

# ==================== THEME-SPECIFIC SLAYDLAR ====================
def add_cover_slide(prs, title, subtitle, theme_key):
    theme = get_theme(theme_key)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, theme["bg_dark"])
    
    # Theme ga mos bezak
    if theme_key in ["astronomik", "fizik"]:
        # Yulduzcha yoki doira
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(1), Inches(4), Inches(4))
        circle.fill.solid(); circle.fill.fore_color.rgb = theme["accent"]
        circle.line.fill.background()
    
    _add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(2), title, 44, RGBColor(255,255,255), bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _add_text(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(1), subtitle, 20, theme["accent"], italic=True, align=PP_ALIGN.CENTER)
    return slide

def add_content_slide(prs, title, bullets, slide_no, total, theme_key):
    theme = get_theme(theme_key)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, theme["bg_light"])
    
    _add_text(slide, Inches(0.8), Inches(0.6), Inches(11.5), Inches(1.2), title, 32, theme["bg_dark"], bold=True)
    
    bullets = [b for b in (bullets or []) if b][:5]
    top = Inches(2.0)
    for i, bullet in enumerate(bullets):
        # Theme ga mos bullet
        if theme_key in ["matematik", "fizik"]:
            _add_text(slide, Inches(0.9), top + i*Inches(1.1), Inches(0.4), Inches(0.8), "◆", 24, theme["accent"], bold=True)
        else:
            _add_text(slide, Inches(0.9), top + i*Inches(1.1), Inches(0.4), Inches(0.8), "•", 28, theme["accent"], bold=True)
        
        _add_text(slide, Inches(1.6), top + i*Inches(1.1), Inches(10.5), Inches(1.0), bullet, 18, theme["text"])
    
    _add_text(slide, Inches(11), Inches(6.8), Inches(2), Inches(0.5), f"{slide_no}/{total}", 11, RGBColor(0x99,0x99,0x99), align=PP_ALIGN.RIGHT)
    return slide

def add_closing_slide(prs, title, subtitle, theme_key):
    theme = get_theme(theme_key)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, theme["bg_dark"])
    _add_text(slide, Inches(1), Inches(2.8), Inches(11.3), Inches(2), title, 40, RGBColor(255,255,255), bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _add_text(slide, Inches(1.5), Inches(4.7), Inches(10.3), Inches(1), subtitle, 18, theme["accent"], align=PP_ALIGN.CENTER)

def create_presentation(slides_data: list, output_path: str, theme: str = "rasmiy"):
    global theme_key
    theme_key = theme  # global for font choice
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    total = len(slides_data)
    for i, item in enumerate(slides_data, 1):
        title = (item.get("title") or "").strip()
        subtitle = (item.get("subtitle") or "").strip()
        bullets = item.get("bullets") or []
        
        if i == 1:
            add_cover_slide(prs, title, subtitle, theme)
        elif i == total and not bullets:
            add_closing_slide(prs, title, subtitle, theme)
        else:
            add_content_slide(prs, title, bullets, i, total, theme)
    
    prs.save(output_path)
    return output_path
    
    
